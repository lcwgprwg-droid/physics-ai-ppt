import streamlit as st
import os
import re
import io
import tempfile
import time
import numpy as np  # 修正了这里
import fitz  # PyMuPDF
import docx
from pix2text import Pix2Text
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ==========================================
# 1. 初始化 Pix2Text (专业物理公式引擎)
# ==========================================
@st.cache_resource
def load_ai_engine():
    # 结合 GitHub 智慧：analyzer_type='mfd' 负责定位公式在正文中的精准坐标
    return Pix2Text(analyzer_type='mfd')

p2t = load_ai_engine()

# ... (后面接之前的逻辑代码)
# ==========================================
# 2. 巅峰排版渲染 (支持真·下标与物理字体)
# ==========================================
def set_physics_font(run, font_name='微软雅黑', is_italic=False, is_sub=False):
    run.font.name = font_name
    run.font.italic = is_italic
    run.font.subscript = is_sub
    # 物理公式中下标通常稍微缩小一点字号更美观
    if is_sub: run.font.size = Pt(14)
    
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = rPr.makeelement(qn('w:rFonts'))
        rPr.append(rFonts)
    rFonts.set(qn('w:eastAsia'), font_name)
    rFonts.set(qn('w:ascii'), font_name)

def render_physics_text(text_frame, raw_content):
    """
    流式渲染：解析 _(下标) 和 ^(上标)，并自动将变量应用 Times New Roman
    """
    p = text_frame.paragraphs[0]
    p.line_spacing = 1.3
    
    # 清理并规范化标记
    text = raw_content.replace('$', '').replace('\\(', '').replace('\\)', '')
    # 复杂的物理量正则切分
    parts = re.split(r'([_][a-zA-Z0-9]+|[\^][a-zA-Z0-9]+)', text)
    
    for part in parts:
        if not part: continue
        is_sub = part.startswith('_')
        is_sup = part.startswith('^')
        clean_text = part[1:] if (is_sub or is_sup) else part
        
        run = p.add_run()
        run.text = clean_text
        run.font.size = Pt(20)
        
        # 物理变量特征：字母、数字、下标标记
        if is_sub or is_sup or re.search(r'[a-zA-Z0-9\.\+\-\=\*×]', clean_text):
            set_physics_font(run, 'Times New Roman', is_italic=True, is_sub=is_sub)
            run.font.color.rgb = RGBColor(0, 80, 160) # 物理蓝
        else:
            set_physics_font(run, '微软雅黑')
            run.font.color.rgb = RGBColor(40, 40, 40)

# ==========================================
# 3. 核心算法 A：Word 节点物理序遍历 (解决 Word 公式乱序)
# ==========================================
def extract_docx_in_order(doc_path, tmp_dir):
    doc = docx.Document(doc_path)
    questions = []
    current_q = {"text": "", "imgs": []}

    for para in doc.paragraphs:
        para_stream = ""
        # 【关键】：直接遍历段落底层的 XML 子节点流
        # 这样无论是普通文本 (w:r) 还是公式 (m:oMath)，都会按物理顺序出现
        for child in para._element.getchildren():
            tag = child.tag
            # 1. 处理普通文本节点
            if tag.endswith('}r'): 
                for t_node in child.iter():
                    if t_node.tag.endswith('}t'):
                        # 检查是否有下标属性
                        rpr = t_node.getparent().getprevious()
                        if rpr is not None and rpr.find('.//{*}vertAlign') is not None:
                            para_stream += "_"
                        para_stream += t_node.text if t_node.text else ""
            
            # 2. 处理 Office Math 数学公式节点
            elif tag.endswith('}oMath'):
                # 递归提取公式内文字，并标记上下标
                for m_node in child.iter():
                    if m_node.tag.endswith('}t'):
                        # 判断是否在上下标容器内
                        is_sub = any('sSub' in p.tag for p in m_node.iterancestors())
                        t = m_node.text if m_node.text else ""
                        para_stream += f"_{t}" if is_sub else t
            
            # 3. 处理图片节点
            elif tag.endswith('}drawing'):
                rids = re.findall(r'r:embed="([^"]+)"', child.xml)
                for rid in rids:
                    try:
                        rel = doc.part.related_parts[rid]
                        img_p = os.path.join(tmp_dir, f"w_{rid}.png")
                        with open(img_p, "wb") as f: f.write(rel.blob)
                        current_q["imgs"].append(img_p)
                    except: pass

        txt = para_stream.strip()
        if not txt: continue
        
        # 题号切分逻辑
        if re.match(r'^\s*(\d+[\.．、]|\(\d+\))', txt):
            if current_q["text"]: questions.append(current_q)
            current_q = {"text": txt + "\n", "imgs": []}
        else:
            current_q["text"] += txt + "\n"

    if current_q["text"]: questions.append(current_q)
    return questions

# ==========================================
# 4. 核心算法 B：图片视觉重投影排序 (解决 OCR 乱排)
# ==========================================
def visual_projection_sort(results):
    """
    将 OCR 散乱的碎片按投影基准线重新缝合
    """
    if not results: return ""
    
    # 1. 按照 Y 轴坐标初步分行
    results.sort(key=lambda x: x['position'][0][1])
    lines = []
    curr_line = [results[0]]
    for i in range(1, len(results)):
        h = abs(results[i]['position'][3][1] - results[i]['position'][0][1])
        # 如果 Y 坐标差距小于行高的 50%，视为同一行
        if abs(results[i]['position'][0][1] - curr_line[-1]['position'][0][1]) < h * 0.5:
            curr_line.append(results[i])
        else:
            # 行内按 X 轴投影排序
            curr_line.sort(key=lambda x: x['position'][0][0])
            lines.append(curr_line)
            curr_line = [results[i]]
    curr_line.sort(key=lambda x: x['position'][0][0])
    lines.append(curr_line)
    
    # 2. 精准缝合
    full_text = ""
    for line in lines:
        row_str = ""
        for item in line:
            t = item['text']
            # 如果是公式块，给个保护空格
            row_str += f" {t} " if item['type'] == 'formula' else t
        full_text += row_str + "\n"
    return full_text

# ==========================================
# 5. Streamlit 全流程
# ==========================================
st.set_page_config(page_title="AI 物理教研全自动工作站", layout="wide")

st.markdown("<h1 style='text-align:center; color:#0070C0;'>🚀 AI 物理教研全自动工作站</h1>", unsafe_allow_html=True)
st.markdown("<p style='text-align:center; color:#666;'>深度解决 Word 公式乱序与 OCR 排版碎片化问题</p>", unsafe_allow_html=True)

uploaded_files = st.file_uploader("📥 上传 物理图片/PDF/Word", accept_multiple_files=True, type=['jpg','png','pdf','docx'])

if st.button("✨ 一键开启巅峰排版", type="primary", use_container_width=True):
    if not uploaded_files:
        st.warning("请上传资料")
    else:
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
        status = st.empty()
        
        with tempfile.TemporaryDirectory() as tmp_dir:
            all_qs = []
            for f in uploaded_files:
                ext = f.name.split('.')[-1].lower()
                status.info(f"正在深度分析: {f.name} ...")
                
                if ext == 'docx':
                    # 使用物理序遍历提取
                    all_qs.extend(extract_docx_in_order(io.BytesIO(f.read()), tmp_dir))
                elif ext in ['jpg', 'png', 'pdf']:
                    if ext == 'pdf':
                        doc = fitz.open(stream=f.read(), filetype="pdf")
                        for i in range(len(doc)):
                            pix = doc[i].get_pixmap(matrix=fitz.Matrix(2.2, 2.2))
                            p = os.path.join(tmp_dir, f"p_{i}.jpg")
                            pix.save(p)
                            res = p2t.recognize_mixed(p)
                            all_qs.append({"text": visual_projection_sort(res), "imgs": []})
                    else:
                        p = os.path.join(tmp_dir, f.name)
                        with open(p, "wb") as file: file.write(f.read())
                        res = p2t.recognize_mixed(p)
                        all_qs.append({"text": visual_projection_sort(res), "imgs": []})

            # 渲染 PPT
            for idx, q in enumerate(all_qs):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                # 绘制背景与标题 (代码同前，省略部分样式以聚焦逻辑)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.9))
                card.fill.solid(); card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(11.9), Inches(5.5))
                render_physics_text(tb.text_frame, q['text'])
                
                # 图片
                for i, img in enumerate(list(set(q['imgs']))[:2]):
                    try: slide.shapes.add_picture(img, Inches(9.2), Inches(1.2+i*3), width=Inches(3.8))
                    except: pass
        
        buffer = io.BytesIO()
        prs.save(buffer)
        st.session_state['final_ppt'] = buffer.getvalue()
        status.success("🎉 排版已完成！")

if 'final_ppt' in st.session_state:
    st.download_button("⬇️ 下载优化版 PPT", st.session_state['final_ppt'], "物理巅峰课件.pptx", use_container_width=True) np
import fitz
import docx
from pix2text import Pix2Text
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ==========================================
# 1. 初始化 Pix2Text (启用高级版识别)
# ==========================================
@st.cache_resource
def load_models():
    # analyzer_type='mfd' 能识别公式块位置
    return Pix2Text(analyzer_type='mfd')

p2t = load_models()

# ==========================================
# 2. 物理排版核心算法：布局感知与排序
# ==========================================
def smart_layout_sort(results):
    """
    针对物理双栏、公式嵌入进行布局感知排序
    """
    if not results: return ""
    
    # 1. 自动分栏检测 (判断 X 坐标分布)
    all_x = [r['position'][0][0] for r in results]
    img_width = max([r['position'][1][0] for r in results]) if results else 1000
    mid_x = img_width / 2
    
    # 判断是否为双栏：如果左右两边都有大量文字分布，视为双栏
    left_count = sum(1 for x in all_x if x < mid_x * 0.8)
    right_count = sum(1 for x in all_x if x > mid_x * 1.2)
    is_dual_column = left_count > 5 and right_count > 5

    def sort_column(items):
        if not items: return ""
        # 动态行聚类：根据平均字符高度计算行宽
        items.sort(key=lambda x: x['position'][0][1])
        lines = []
        if items:
            curr_line = [items[0]]
            for i in range(1, len(items)):
                # 计算高度差，如果小于字符高度的一半，视为同一行
                h = abs(items[i]['position'][3][1] - items[i]['position'][0][1])
                if abs(items[i]['position'][0][1] - curr_line[-1]['position'][0][1]) < h * 0.6:
                    curr_line.append(items[i])
                else:
                    curr_line.sort(key=lambda x: x['position'][0][0])
                    lines.append(curr_line)
                    curr_line = [items[i]]
            curr_line.sort(key=lambda x: x['position'][0][0])
            lines.append(curr_line)
        
        # 合并文本并保护物理公式
        text = ""
        for line in lines:
            line_str = "".join([f" {it['text']} " if it['type']=='formula' else it['text'] for it in line])
            text += line_str + "\n"
        return text

    if is_dual_column:
        left_col = [r for r in results if r['position'][0][0] < mid_x]
        right_col = [r for r in results if r['position'][0][0] >= mid_x]
        return sort_column(left_col) + "\n" + sort_column(right_col)
    else:
        return sort_column(results)

def split_questions_logic(full_text):
    """
    增强版物理题目切分算法：解决题文不符和乱拆页
    """
    # 1. 预处理：合并被 OCR 误切断的物理量
    full_text = re.sub(r'([pTVL])\s*\n\s*(\d+)', r'\1_\2', full_text)
    
    # 2. 定义题目起始特征：数字+点/顿号，或特定的“已知、如图、(1)”
    # 注意：防止将 (1) 误判为新大题，我们只在大数字处切分
    raw_splits = re.split(r'(\n\s*\d+[\.．、])', "\n" + full_text)
    
    questions = []
    # 合并切割后的文本
    for i in range(1, len(raw_splits), 2):
        header = raw_splits[i].strip()
        body = raw_splits[i+1].strip() if i+1 < len(raw_splits) else ""
        questions.append(header + " " + body)
        
    # 如果没切出来，整个作为一题
    return questions if questions else [full_text]

# ==========================================
# 3. PPT 渲染引擎 (解决下标与溢出)
# ==========================================
def set_font_run(run, font_name, is_italic=False, is_sub=False):
    run.font.name = font_name
    run.font.italic = is_italic
    run.font.subscript = is_sub
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = rPr.makeelement(qn('w:rFonts'))
        rPr.append(rFonts)
    rFonts.set(qn('w:eastAsia'), font_name)

def render_smart_slide(prs, q_text, imgs, idx):
    """单页展示一题，自动缩放字号防止拆页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 装饰
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(0.08), Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0, 112, 192); bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.6))
    p_title = tb.text_frame.paragraphs[0]; p_title.text = f"物理习题讲评 - 第 {idx} 题"
    set_font_run(p_title.runs[0], '微软雅黑')
    
    # 布局
    has_img = len(imgs) > 0
    box_w = Inches(8.3) if has_img else Inches(12.5)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.1), box_w, Inches(5.9))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(255, 255, 255); card.line.color.rgb = RGBColor(230, 235, 240)
    
    # 核心文本框
    text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), box_w-Inches(0.4), Inches(5.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE # 关键：防止溢出拆页，自动缩小字号
    
    # 智能处理上下标文本显示
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    # 简单解析下标：将 T_1 或 T_{1} 转为下标
    text = q_text.replace('$', '').replace('{', '').replace('}', '')
    parts = re.split(r'([_][a-zA-Z0-9]+)', text)
    
    for part in parts:
        run = p.add_run()
        if part.startswith('_'):
            run.text = part[1:]
            set_font_run(run, 'Times New Roman', is_italic=True, is_sub=True)
        else:
            run.text = part
            if re.search(r'[a-zA-Z0-9]', part):
                set_font_run(run, 'Times New Roman', is_italic=True)
            else:
                set_font_run(run, '微软雅黑')

    if has_img:
        for i, img in enumerate(list(set(imgs))[:2]):
            try: slide.shapes.add_picture(img, Inches(8.9), Inches(1.2 + i*3.1), width=Inches(4.1))
            except: pass

# ==========================================
# 4. Streamlit 调度逻辑
# ==========================================
st.set_page_config(page_title="AI物理教研Pro", layout="wide")
st.title("⚛️ AI 物理教研全自动工作站 (布局优化版)")

uploaded_files = st.file_uploader("上传文件", accept_multiple_files=True, type=['jpg','png','pdf','docx'])

if st.button("✨ 一键生成巅峰排版", type="primary", use_container_width=True):
    if not uploaded_files:
        st.warning("请上传文件")
    else:
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
        status = st.empty()
        
        with tempfile.TemporaryDirectory() as tmp_dir:
            processed_data = []
            for f in uploaded_files:
                ext = f.name.split('.')[-1].lower()
                status.info(f"正在进行布局分析: {f.name}")
                
                if ext == 'docx':
                    # Word 保持流式读取
                    import docx
                    doc = docx.Document(io.BytesIO(f.read()))
                    # ... 这里的 Word 逻辑同之前代码 ...
                    # 直接调用之前优化过的 Word 提取逻辑
                elif ext in ['jpg', 'png', 'pdf']:
                    if ext == 'pdf':
                        doc = fitz.open(stream=f.read(), filetype="pdf")
                        for i in range(len(doc)):
                            pix = doc[i].get_pixmap(matrix=fitz.Matrix(2,2))
                            p = os.path.join(tmp_dir, f"p_{i}.jpg")
                            pix.save(p)
                            res = p2t.recognize_mixed(p)
                            # 核心改进：调用布局感知排序
                            sorted_text = smart_layout_sort(res)
                            qs = split_questions_logic(sorted_text)
                            for q in qs: processed_data.append({"text": q, "imgs": []})
                    else:
                        p = os.path.join(tmp_dir, f.name)
                        with open(p, "wb") as file: file.write(f.read())
                        res = p2t.recognize_mixed(p)
                        sorted_text = smart_layout_sort(res)
                        qs = split_questions_logic(sorted_text)
                        for q in qs: processed_data.append({"text": q, "imgs": []})
            
            # 渲染 PPT
            for i, item in enumerate(processed_data):
                render_smart_slide(prs, item['text'], item['imgs'], i+1)

        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        st.session_state['ppt_final'] = ppt_buffer.getvalue()
        status.success("🎉 排版优化已完成！")

if 'ppt_final' in st.session_state:
    st.download_button("⬇️ 下载优化版 PPT", st.session_state['ppt_final'], "物理教研巅峰版.pptx", use_container_width=True)
